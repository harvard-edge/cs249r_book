#!/usr/bin/env python3
"""Convert PDF slide decks to PowerPoint (PPTX) using high-resolution images.

Each PDF page is rendered at 300 DPI via pdftoppm (poppler) and placed as a
full-bleed image on a 16:9 PowerPoint slide. The result is visually identical
to the PDF — suitable for presenting in PowerPoint/Keynote with annotations,
presenter view, etc. Slides are image-based (not editable text).

Requirements:
    - pdftoppm (poppler-utils): apt install poppler-utils / brew install poppler
    - python-pptx:             pip install python-pptx
    - Pillow:                  pip install Pillow

Usage:
    python3 pdf2pptx.py input.pdf                    # -> input.pptx
    python3 pdf2pptx.py input.pdf -o output.pptx     # explicit output
    python3 pdf2pptx.py _build/vol1/*.pdf             # batch convert
    python3 pdf2pptx.py _build/vol1/*.pdf --dpi 200   # lower DPI (smaller files)
"""

import argparse
import glob
import os
import shutil
import subprocess
import sys
import tempfile

from pathlib import Path

def check_dependencies():
    if not shutil.which("pdftoppm"):
        print("ERROR: pdftoppm not found. Install poppler:", file=sys.stderr)
        print("  macOS:  brew install poppler", file=sys.stderr)
        print("  Ubuntu: sudo apt install poppler-utils", file=sys.stderr)
        sys.exit(1)
    try:
        from pptx import Presentation  # noqa: F401
    except ImportError:
        print("ERROR: python-pptx not found. Install:", file=sys.stderr)
        print("  pip install python-pptx Pillow", file=sys.stderr)
        sys.exit(1)


def convert_pdf_to_pptx(pdf_path: str, pptx_path: str, dpi: int = 300) -> dict:
    """Convert a single PDF to PPTX. Returns stats dict."""
    from pptx import Presentation
    from pptx.util import Inches, Emu

    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run(
            ["pdftoppm", "-png", "-r", str(dpi), pdf_path, os.path.join(tmp, "slide")],
            check=True,
            capture_output=True,
        )
        images = sorted(glob.glob(os.path.join(tmp, "slide-*.png")))
        if not images:
            return {"ok": False, "error": "pdftoppm produced no images"}

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for img_path in images:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            slide.shapes.add_picture(
                img_path, Emu(0), Emu(0), prs.slide_width, prs.slide_height
            )

        os.makedirs(os.path.dirname(pptx_path) or ".", exist_ok=True)
        prs.save(pptx_path)

        size_mb = os.path.getsize(pptx_path) / (1024 * 1024)
        return {"ok": True, "slides": len(images), "size_mb": size_mb}


def main():
    parser = argparse.ArgumentParser(
        description="Convert PDF slide decks to PowerPoint (image-based)"
    )
    parser.add_argument("pdfs", nargs="+", help="PDF file(s) to convert")
    parser.add_argument("-o", "--output", help="Output path (single file only)")
    parser.add_argument("--dpi", type=int, default=300, help="Render DPI (default: 300)")
    parser.add_argument("--output-dir", help="Output directory (batch mode)")
    args = parser.parse_args()

    check_dependencies()

    if args.output and len(args.pdfs) > 1:
        print("ERROR: -o/--output only works with a single input file", file=sys.stderr)
        sys.exit(1)

    total = 0
    failed = 0

    for pdf in args.pdfs:
        if not os.path.isfile(pdf):
            print(f"  SKIP {pdf} (not found)")
            continue

        if args.output:
            pptx = args.output
        elif args.output_dir:
            pptx = os.path.join(args.output_dir, Path(pdf).stem + ".pptx")
        else:
            pptx = str(Path(pdf).with_suffix(".pptx"))

        print(f"  Converting {pdf}...")
        result = convert_pdf_to_pptx(pdf, pptx, dpi=args.dpi)

        if result["ok"]:
            print(f"    -> {pptx} ({result['slides']} slides, {result['size_mb']:.1f} MB)")
            total += 1
        else:
            print(f"    FAILED: {result['error']}")
            failed += 1

    print(f"\nDone: {total} converted, {failed} failed")
    if failed:
        sys.exit(1)


if __name__ == "__main__":
    main()
